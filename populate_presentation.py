"""Populate IBM capstone PowerPoint template and export submission PDF."""

from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "ds-capstone-template-coursera.pptx"
OUTPUT_PPTX = ROOT / "Data Science Capstone Project Report.pptx"
OUTPUT_PDF = ROOT / "Data Science Capstone Project Report.pdf"
IMAGES = ROOT / "images"

AUTHOR = "Vikas Parmar"
GITHUB = "https://github.com/vikas-parmar/applied-data-science-capstone-ibm"
TODAY = date.today().strftime("%B %d, %Y")

NOTEBOOK_LINKS = {
    "api": f"{GITHUB}/blob/main/Data%20Collection%20API.ipynb",
    "scrape": f"{GITHUB}/blob/main/Data%20Collection%20with%20Web%20Scraping.ipynb",
    "wrangle": f"{GITHUB}/blob/main/Data%20Wrangling.ipynb",
    "viz": f"{GITHUB}/blob/main/EDA%20with%20Data%20Visualization.ipynb",
    "sql": f"{GITHUB}/blob/main/EDA.ipynb",
    "folium": f"{GITHUB}/blob/main/Interactive%20Visual%20Analytics%20with%20Folium%20lab.ipynb",
    "dash": f"{GITHUB}/blob/main/spacex_dash_app.py",
    "ml": f"{GITHUB}/blob/main/Machine%20Learning%20Prediction.ipynb",
}


def replace_slide_text(slide, replacements: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        updated = text
        for old, new in replacements.items():
            updated = updated.replace(old, new)
        if updated != text:
            shape.text_frame.text = updated


def set_body_text(slide, text: str) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(16)
            return
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.text = text


def add_image(slide, image_path: Path, left=1.0, top=1.8, width=8.5) -> None:
    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), width=Inches(width)
        )


def populate_presentation() -> Presentation:
    prs = Presentation(str(TEMPLATE))

    replace_slide_text(
        prs.slides[0],
        {"<Name>": AUTHOR, "<Date>": TODAY},
    )

    executive = (
        "This project predicts SpaceX Falcon 9 first-stage landing success using data "
        "science methods across the full workflow: API data collection, web scraping, "
        "data wrangling, SQL and visual EDA, Folium mapping, Plotly Dash analytics, and "
        "machine learning classification.\n\n"
        "Key result: Decision Tree achieved the best cross-validation score among four "
        "models, and EDA shows launch site, orbit type, payload mass, and flight "
        "experience correlate with landing outcomes."
    )
    set_body_text(prs.slides[2], executive)

    intro = (
        "SpaceX advertises Falcon 9 launches at about $62M versus $165M+ from other "
        "providers. Reusing the first stage drives much of that savings.\n\n"
        "Problem: For a given launch (payload, orbit, site, booster version), will the "
        "Falcon 9 first stage land successfully?\n\n"
        "Audience: A competing aerospace startup estimating launch bid costs against SpaceX."
    )
    set_body_text(prs.slides[3], intro)

    methodology = (
        "Data collection: SpaceX REST API and Wikipedia scraping\n"
        "Data wrangling: Missing values, outcome mapping, Class label\n"
        "EDA: SQL queries + Matplotlib/Seaborn charts\n"
        "Interactive analytics: Folium maps + Plotly Dash dashboard\n"
        "Predictive analysis: StandardScaler, GridSearchCV, 4 classifiers"
    )
    set_body_text(prs.slides[5], methodology)

    slide8 = (
        "Flowchart:\n"
        "SpaceX API -> requests.get() -> json_normalize() -> Filter Falcon 9 -> "
        "Fill missing PayloadMass -> dataset_part_1.csv\n\n"
        f"GitHub: {NOTEBOOK_LINKS['api']}"
    )
    set_body_text(prs.slides[7], slide8)

    slide9 = (
        "Flowchart:\n"
        "Wikipedia page -> requests.get() -> BeautifulSoup -> parse HTML table -> "
        "Pandas DataFrame -> spacex_web_scraped.csv\n\n"
        f"GitHub: {NOTEBOOK_LINKS['scrape']}"
    )
    set_body_text(prs.slides[8], slide9)

    slide10 = (
        "Flowchart:\n"
        "dataset_part_1.csv -> EDA summaries -> map Outcome to Class (0/1) -> "
        "dataset_part_2.csv + spacex_launch_dash.csv\n\n"
        f"GitHub: {NOTEBOOK_LINKS['wrangle']}"
    )
    set_body_text(prs.slides[9], slide10)

    slide11 = (
        "Charts plotted:\n"
        "- Flight Number vs Launch Site\n"
        "- Payload vs Launch Site\n"
        "- Success Rate vs Orbit Type\n"
        "- Flight Number vs Orbit Type\n"
        "- Payload vs Orbit Type\n"
        "- Launch Success Yearly Trend\n\n"
        f"GitHub: {NOTEBOOK_LINKS['viz']}"
    )
    set_body_text(prs.slides[10], slide11)

    slide12 = (
        "SQL queries performed:\n"
        "SELECT DISTINCT launch sites; CCA launch sites; NASA CRS payload sum; "
        "F9 v1.1 average payload; first ground landing date; drone ship payload range; "
        "success/failure counts; max payload boosters; 2015 failures; ranked outcomes.\n\n"
        f"GitHub: {NOTEBOOK_LINKS['sql']}"
    )
    set_body_text(prs.slides[11], slide12)

    slide13 = (
        "Folium map objects:\n"
        "- Markers for all launch sites\n"
        "- Color-coded success/failure launch markers\n"
        "- Proximity lines to coastline, highway, railway, and city\n\n"
        f"GitHub: {NOTEBOOK_LINKS['folium']}"
    )
    set_body_text(prs.slides[12], slide13)

    slide14 = (
        "Plotly Dash dashboard:\n"
        "- Launch site dropdown\n"
        "- Success pie chart\n"
        "- Payload range slider\n"
        "- Payload vs launch outcome scatter plot\n\n"
        f"GitHub: {NOTEBOOK_LINKS['dash']}"
    )
    set_body_text(prs.slides[13], slide14)

    slide15 = (
        "Predictive workflow:\n"
        "Merge datasets -> StandardScaler -> train/test split -> GridSearchCV for "
        "Logistic Regression, SVM, Decision Tree, KNN -> evaluate accuracy and confusion matrix\n\n"
        f"GitHub: {NOTEBOOK_LINKS['ml']}"
    )
    set_body_text(prs.slides[14], slide15)

    eda_images = [
        (17, "eda_01_flight_vs_site.png", "Flight number increases over time at each launch site."),
        (18, "eda_02_payload_vs_site.png", "Payload distribution varies significantly by launch site."),
        (19, "eda_03_success_vs_orbit.png", "LEO and ISS orbits show higher success rates than GTO."),
        (20, "eda_04_flight_vs_orbit.png", "Later flights show more diverse orbit assignments."),
        (21, "eda_05_payload_vs_orbit.png", "Heavier payloads appear more often in GTO missions."),
        (22, "eda_06_yearly_trend.png", "Landing success rate improved sharply after 2014."),
    ]
    for idx, img, note in eda_images:
        add_image(prs.slides[idx], IMAGES / img)
        set_body_text(prs.slides[idx], note)

    sql_images = [
        (23, "sql_01_launch_sites.png", "Four unique launch sites identified."),
        (24, "sql_02_cca_sites.png", "CCAFS launch sites begin with CCA prefix."),
        (25, "sql_03_nasa_payload.png", "Total NASA CRS payload mass calculated."),
        (26, "sql_04_avg_f9v11.png", "Average payload for F9 v1.1 boosters."),
        (27, "sql_05_first_ground_success.png", "First successful RTLS landing date."),
        (28, "sql_06_drone_ship_payload.png", "Boosters with drone ship landing and mid-range payload."),
        (29, "sql_07_mission_outcomes.png", "Successful vs failed mission counts."),
        (30, "sql_08_max_payload.png", "Boosters carrying maximum payload mass."),
        (31, "sql_09_failed_2015.png", "Failed drone ship landings in 2015."),
        (32, "sql_10_rank_outcomes.png", "Ranked landing outcomes between 2010 and 2017."),
    ]
    for idx, img, note in sql_images:
        add_image(prs.slides[idx], IMAGES / img, top=1.6, width=7.5)
        set_body_text(prs.slides[idx], note)

    folium_images = [
        (34, "folium_1.png", "All four SpaceX launch sites marked on the map."),
        (35, "folium_2.png", "Green markers = successful landings; red = failures."),
        (36, "folium_3.png", "Proximity analysis shows distance to key infrastructure."),
    ]
    for idx, img, note in folium_images:
        add_image(prs.slides[idx], IMAGES / img)
        set_body_text(prs.slides[idx], note)

    dash_images = [
        (38, "dash_01_all_sites_pie.png", "Success share by launch site across all locations."),
        (39, "dash_02_site_pie.png", "Site with highest success ratio shown in pie chart."),
        (40, "dash_03_payload_scatter.png", "Payload vs outcome for selected payload range."),
    ]
    for idx, img, note in dash_images:
        add_image(prs.slides[idx], IMAGES / img)
        set_body_text(prs.slides[idx], note)

    summary = json.loads((IMAGES / "ml_summary.json").read_text(encoding="utf-8"))
    best = summary["best_model"]
    add_image(prs.slides[42], IMAGES / "ml_accuracy_bar.png")
    set_body_text(
        prs.slides[42],
        f"Model comparison shows test accuracy for all four classifiers. "
        f"{best} has the highest GridSearchCV score.",
    )

    add_image(prs.slides[43], IMAGES / "ml_confusion_matrix.png", width=5.5)
    set_body_text(
        prs.slides[43],
        f"Confusion matrix for {best}. False positives would overestimate reusable "
        "first-stage availability and understate launch cost.",
    )

    conclusions = (
        "1. Landing success improved significantly after 2014 as SpaceX gained experience.\n"
        "2. Orbit type and payload mass strongly influence landing outcome.\n"
        "3. Launch site selection affects both success rate and operational risk.\n"
        "4. Decision Tree is the recommended model for predicting first-stage landing success."
    )
    set_body_text(prs.slides[44], conclusions)

    appendix = (
        f"Author: {AUTHOR}\n"
        f"GitHub Repository: {GITHUB}\n"
        "Additional assets: notebooks, SQL query outputs, charts, and CSV datasets."
    )
    set_body_text(prs.slides[45], appendix)

    try:
        for slide in prs.slides:
            watermark = slide.shapes.add_textbox(
                Inches(9.5), Inches(7.0), Inches(3.5), Inches(0.4)
            )
            watermark.text_frame.text = AUTHOR
            watermark.text_frame.paragraphs[0].font.size = Pt(10)
            watermark.text_frame.paragraphs[0].font.italic = True
    except Exception:
        pass

    prs.save(str(OUTPUT_PPTX))
    return prs


def export_pdf() -> None:
    pptx_path = str(OUTPUT_PPTX.resolve())
    pdf_path = str(OUTPUT_PDF.resolve())

    try:
        import comtypes.client

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        return
    except Exception:
        pass

    try:
        import win32com.client

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        return
    except Exception:
        pass

    raise RuntimeError(
        "Could not export PDF automatically. Open "
        f"{OUTPUT_PPTX.name} in PowerPoint and export as PDF manually."
    )


def main() -> None:
    print("Populating presentation...")
    populate_presentation()
    print(f"Saved {OUTPUT_PPTX.name}")
    print("Exporting PDF...")
    try:
        export_pdf()
        print(f"Saved {OUTPUT_PDF.name}")
    except RuntimeError as exc:
        print(exc)


if __name__ == "__main__":
    main()
